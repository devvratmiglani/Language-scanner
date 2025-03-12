import PyPDF2
import docx
import pandas as pd
from pptx import Presentation

def extract_text(file_path, file_type):
    text = ""
    try:
        if file_type == "pdf":
            with open(file_path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    text += page.extract_text() or ""
        elif file_type == "docx":
            doc = docx.Document(file_path)
            text = "\n".join([para.text for para in doc.paragraphs])
        elif file_type == "xlsx":
            df = pd.read_excel(file_path)
            text = df.to_string()
        elif file_type == "csv":
            df = pd.read_csv(file_path)
            text = df.to_string()
        elif file_type == "pptx":
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        elif file_type == "txt":
            with open(file_path, "r", encoding="utf-8") as f:
                text = f.read()
        else:
            with open(file_path, "rb") as f:
                text = f.read().decode(errors='ignore')
    except Exception as e:
        return f"Error extracting text: {str(e)}"
    return text.strip()
